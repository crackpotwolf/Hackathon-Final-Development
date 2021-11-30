from pptx import Presentation, table, slide


class PptxLoader:
  """Класс обработки файлов формата *.pptx
  """
  def __init__(self):
    """Метод инициализации объекта класса
    """
    self.data = {}
    
  def __read_file__(self, filepath: str):
    """Метод считывания данных из файла

    Args:
        filepath (str): [путь к файлу]
    """

    def get_data_from_table(table: table):
      """Метод считывания данных из таблицы

      Args:
          table (table): [текущая таблицы]

      Returns:
          [List]: [лист словарей, содержащие данные по ячейкам в формате название столбца: значение ячейки]
      """
      # метод парсинга таблицы
      data = []
      keys = None
      # проход по всем строкам таблицы
      for i, row in enumerate(table.rows):
        # забрать текст из ячейки
        text = (cell.text for cell in row.cells)
        # если строка первая (0), то  задать заголовки столбцов
        if i == 0:
            keys = tuple(text)
            continue
        # запаковывание данных ячейки
        row_data = dict(zip(keys, text))
        data.append(row_data)
      return data
    
    def get_shapes_data(current_slide: slide, num_slide: int):
      """Парсинг объектов на слайде

      Args:
          current_slide (slide): [текущий слайд]
          num_slide (int): [номер слайд]
      """
      # получить объекты на слайде
      slide_shapes = current_slide.shapes
      # номер объекта - титула
      title_shape_id = current_slide.title.shape_id
      # шаблон названия ключа
      key_name = "slide_%s_"%str(num_slide)
      
      # определение типа объекта и считывание данных
      table_count = 0
      text_count = 0
      # обработка объектов
      for current_shape in slide_shapes:
        # если объект - таблица
        if current_shape.has_table:
          table_count += 1
          self.data[key_name+"text_%s"%str(table_count)] = get_data_from_table(table=current_shape.table)
        
        # если объект - текстовый элемент
        elif current_shape.has_text_frame:
          # если титул, то создать ключ с пометкой титула
          
          if current_shape.shape_id == title_shape_id:
            self.data[key_name+"text_title"] = current_shape.text
          
          else:
            text_count += 1
            self.data[key_name+"text_%s"%str(text_count)] = current_shape.text
    
    
    def get_data_from_slides(presentation: Presentation):
      """Метод считывания данных со слайдов

      Args:
          presentation (Presentation): [текущая презентация]
      """
      # получить все слайды
      slides = presentation.slides
      # получить данные со слайдов
      for num_slide, current_slide in enumerate(slides):
        get_shapes_data(current_slide=current_slide, num_slide=num_slide+1)

    # открыть презентацию
    presentation = Presentation(pptx=filepath)
    # получить данные из презентации
    get_data_from_slides(presentation=presentation)
